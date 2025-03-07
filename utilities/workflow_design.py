# -*- coding: utf-8 -*-
import os
from graphviz import Digraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# ============================================================
# FunÃ§Ãµes Auxiliares para cores (Graphviz e PPTX)
# ============================================================
def hex_to_rgb(hex_color: str) -> tuple:
    """Converte cor hexadecimal para uma tupla RGB."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def is_dark_color(hex_color: str, threshold: int = 130) -> bool:
    """Determina se uma cor Ã© considerada escura com base na luminÃ¢ncia."""
    r, g, b = hex_to_rgb(hex_color)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < threshold

# ============================================================
# PARTE 1 â€“ FLUXOGRAMA COM GRAPHVIZ (PNG, SVG e PDF com DPI 900)
# ============================================================
def create_flowchart():
    # Usando orientaÃ§Ã£o vertical (TB)
    dot = Digraph(format='png')
    # Ajustado para dimensÃµes maiores: adequado para publicaÃ§Ã£o na NAR
    dot.attr(rankdir='TB', size='16,10!', dpi='900')
    # Aumenta fonte e dimensÃµes dos nÃ³s
    dot.attr('node', fontname='Arial', fontsize='20', width='2.5', height='1.5')
    dot.attr('edge', fontname='Arial', fontsize='16')
    
    # Definindo um dicionÃ¡rio de cores em escala de cinza suaves para os nÃ³s
    gray_colors = {
        'A': '#F2F2F2',
        'X': '#F8F8F8',  # ExtraÃ§Ã£o FAAL para Treinamento
        'B': '#E6E6E6',
        'C': '#D9D9D9',
        'D': '#CCCCCC',
        'E': '#C0C0C0',
        'F1': '#C0C0C0',  # Mesma cor de E
        'F': '#B3B3B3',   # Gerar vetores Word2Vec e agregÃ¡-los
        'F2': '#B3B3B3',  # Salvar o modelo global Word2Vec
        'G': '#A6A6A6',   # Standardize Training Embeddings (StandardScaler)
        'H': '#999999',   # SerÃ¡ usado para Oversampling (Random oversampling and SMOTE)
        'H2': '#949494',  # Novo nÃ³: Grid Search and Cross validation
        'I': '#8C8C8C',   # Calibrate Model (Isotonic Calibration)
        'J': '#808080',   # Evaluate Model (ROC, F1, PR AUC)
        'Y': '#7A7A7A',   # ExtraÃ§Ã£o FAAL para PrediÃ§Ã£o
        'K': '#737373',
        'L': '#666666',
        'M': '#5A5A5A',
        'N': '#4D4D4D',
        'O': '#404040',
        'P': '#404040',   # Import Random Forest Classifier and RF Scaler
        'Z': '#505050',   # Sum Probabilities and Normalize
        'T': '#414141',   # Get Prediction Confidence
        'Q': '#333333',   # Visualize Results (Scatter Plot, UMAP, Learning Curve)
        'R': '#262626'    # Output Results (Save Files, Download)
    }
    
    # FunÃ§Ã£o para definir atributos do nÃ³, incluindo fontcolor se a cor for escura
    def add_node(node_id, label, shape, fillcolor):
        attrs = {'shape': shape, 'style': 'filled', 'fillcolor': fillcolor}
        if is_dark_color(fillcolor):
            attrs['fontcolor'] = 'white'
        dot.node(node_id, label, **attrs)
    
    # Usando etiquetas HTML para negrito; escapando "&" como "&amp;"
    # Treinamento
    add_node('A', '<<B>Start: Configure Streamlit &amp; Environment</B>>', 'ellipse', gray_colors['A'])
    add_node('X', '<<B>Extract FAAL Domain\n from Training FASTA</B>>', 'box', gray_colors['X'])
    add_node('B', '<<B>Load Training Data:\n FASTA and Table</B>>', 'parallelogram', gray_colors['B'])
    add_node('C', '<<B>Check Training: Sequence Alignment</B>>', 'box', gray_colors['C'])
    add_node('D', '<<B>Realign with MAFFT</B>>', 'diamond', gray_colors['D'])
    add_node('E', '<<B>Use Aligned Training Data</B>>', 'box', gray_colors['E'])
    add_node('F1', '<<B>Breaking alignment into k-mers</B>>', 'box', gray_colors['F1'])
    add_node('F', '<<B>Generate Word2Vec vectors for each k-mer, aggregate them by mean to form a single sentence vector per sequence</B>>', 'box', gray_colors['F'])
    add_node('F2', '<<B>Save the global Word2Vec model</B>>', 'box', gray_colors['F2'])
    add_node('G', '<<B>Standardize Training Embeddings\n (StandardScaler)</B>>', 'box', gray_colors['G'])
    
    # Dividindo a etapa de treinamento apÃ³s G em dois passos:
    add_node('H', '<<B>Oversampling (Random oversampling and SMOTE)</B>>', 'box', gray_colors['H'])
    add_node('H2', '<<B>Grid Search and Cross validation</B>>', 'box', gray_colors['H2'])
    add_node('I', '<<B>Calibrate Model (Isotonic Calibration)</B>>', 'box', gray_colors['I'])
    add_node('J', '<<B>Evaluate Model (ROC, F1, PR AUC)</B>>', 'box', gray_colors['J'])
    
    # PrediÃ§Ã£o
    add_node('Y', '<<B>Extract FAAL Domain from Prediction FASTA</B>>', 'box', gray_colors['Y'])
    add_node('K', '<<B>Load Prediction Data (FASTA)</B>>', 'parallelogram', gray_colors['K'])
    add_node('L', '<<B>Extract k-mers from new sequences</B>>', 'box', gray_colors['L'])
    add_node('M', '<<B>Extract k-mer embeddings using a pre-trained global Word2Vec model and merge them into a single sentence vector</B>>', 'box', gray_colors['M'])
    add_node('N', '<<B>Standardize New Sequence Embeddings\n(Using Training Scaler)</B>>', 'box', gray_colors['N'])
    add_node('O', '<<B>Import Random Forest Classifier and RF Scaler</B>>', 'box', gray_colors['O'])
    # Corrigido: "&" substituÃ­do por "&amp;"
    add_node('P', '<<B>Predict on New Sequences &amp; Get Class Rankings</B>>', 'box', gray_colors['P'])
    add_node('Z', '<<B>Sum Probabilities and Normalize</B>>', 'box', gray_colors['Z'])
    add_node('T', '<<B>Get Prediction Confidence</B>>', 'box', gray_colors['T'])
    add_node('Q', '<<B>Visualize Results (Scatter Plot, UMAP, Learning Curve)</B>>', 'box', gray_colors['Q'])
    add_node('R', '<<B>Output Results (Save Files, Download)</B>>', 'ellipse', gray_colors['R'])
    
    # Dividindo o fluxo em 4 clusters (painÃ©is)
    with dot.subgraph(name='cluster_A') as cA:
        cA.attr(label='A) Neural Network: Protein Embedding Phase', fontsize='24', fontname='Arial Bold', fontcolor='black', style='dashed', color='black')
        for n in ['A', 'X', 'B', 'C', 'D', 'E', 'F1', 'F', 'F2', 'G']:
            cA.node(n)
    with dot.subgraph(name='cluster_B') as cB:
        cB.attr(label='B) Oversampling Phase', fontsize='24', fontname='Arial Bold', fontcolor='black', style='dashed', color='black')
        for n in ['H', 'H2']:
            cB.node(n)
    with dot.subgraph(name='cluster_C') as cC:
        cC.attr(label='C) Training and Optimization of Random Forest Phase', fontsize='24', fontname='Arial Bold', fontcolor='black', style='dashed', color='black')
        for n in ['I', 'J']:
            cC.node(n)
    with dot.subgraph(name='cluster_D') as cD:
        cD.attr(label='D) Classification of New Sequences Phase â€“ Transfer Learning', fontsize='24', fontname='Arial Bold', fontcolor='black', style='dashed', color='black')
        for n in ['Y', 'K', 'L', 'M', 'N', 'O', 'P', 'Z', 'T', 'Q', 'R']:
            cD.node(n)
    
    # Define as arestas:
    # Treinamento:
    dot.edge('A', 'X')
    dot.edge('X', 'B')
    dot.edge('B', 'C')
    dot.edge('C', 'D', label='Not Aligned')
    dot.edge('C', 'E', label='Aligned')
    dot.edge('D', 'E')
    dot.edge('E', 'F1')
    dot.edge('F1', 'F')
    dot.edge('F', 'F2')
    dot.edge('F2', 'G')
    dot.edge('G', 'H')
    dot.edge('H', 'H2')
    dot.edge('H2', 'I')
    dot.edge('I', 'J')
    # TransiÃ§Ã£o para prediÃ§Ã£o:
    dot.edge('J', 'Y')
    # PrediÃ§Ã£o:
    dot.edge('Y', 'K')
    dot.edge('K', 'L')
    dot.edge('L', 'M')
    dot.edge('M', 'N')
    dot.edge('N', 'O')
    dot.edge('O', 'P')
    dot.edge('P', 'Z')
    dot.edge('Z', 'T')
    dot.edge('T', 'Q')
    dot.edge('Q', 'R')
    
    return dot

def generate_flowchart_images(filename_base="workflow"):
    flowchart = create_flowchart()
    # Gera PNG, SVG e PDF com DPI 900
    png_path = flowchart.render(filename_base, format='png', view=False, cleanup=True)
    svg_path = flowchart.render(filename_base, format='svg', view=False, cleanup=True)
    pdf_path = flowchart.render(filename_base, format='pdf', view=False, cleanup=True)
    print(f"PNG gerado em: {png_path}")
    print(f"SVG gerado em: {svg_path}")
    print(f"PDF gerado em: {pdf_path}")

# ============================================================
# PARTE 2 â€“ FLUXOGRAMA COMO PPTX (FORMAS NATIVAS â€“ ORGANIZADOS VERTICALMENTE)
# ============================================================
def map_shape(shape_str: str):
    """
    Mapeia o tipo de forma para a forma nativa do PowerPoint usando MSO_AUTO_SHAPE_TYPE.
    """
    if shape_str == 'ellipse':
        return MSO_AUTO_SHAPE_TYPE.OVAL
    elif shape_str == 'parallelogram':
        return MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM
    elif shape_str == 'diamond':
        return MSO_AUTO_SHAPE_TYPE.DIAMOND
    elif shape_str == 'box':
        return MSO_AUTO_SHAPE_TYPE.RECTANGLE
    else:
        return MSO_AUTO_SHAPE_TYPE.RECTANGLE

def ppt_font_color(fill_hex: str) -> str:
    """
    Retorna 'FFFFFF' se a cor de preenchimento for escura, senÃ£o '000000'.
    """
    return "FFFFFF" if is_dark_color(fill_hex) else "000000"

# Atualizando os nÃ³s para incluir as novas etapas e ajustando posiÃ§Ãµes (em polegadas)
ppt_nodes = {
    # Painel A: Neural Network: Protein Embedding Phase (nÃ³s A, X, B, C, D, E, F1, F, F2, G)
    'A': {'text': 'Start: Configure\nStreamlit &amp; Environment', 'shape': 'ellipse',       'fillcolor': '#F2F2F2', 'x': 1,   'y': 0.5,  'w': 7, 'h': 2.5},
    'X': {'text': 'Extract FAAL Domain from Training FASTA', 'shape': 'box', 'fillcolor': '#F8F8F8', 'x': 1,   'y': 3.0,  'w': 7, 'h': 2.5},
    'B': {'text': 'Load Training Data:\nFASTA and Table', 'shape': 'parallelogram', 'fillcolor': '#E6E6E6', 'x': 1,   'y': 5.5,  'w': 7, 'h': 2.5},
    'C': {'text': 'Check Training Sequence Alignment', 'shape': 'box', 'fillcolor': '#D9D9D9', 'x': 1,   'y': 8.0,  'w': 7, 'h': 2.5},
    'D': {'text': 'Realign with MAFFT', 'shape': 'diamond', 'fillcolor': '#CCCCCC', 'x': 1,   'y': 10.5, 'w': 7, 'h': 2.5},
    'E': {'text': 'Use Aligned Training Data', 'shape': 'box', 'fillcolor': '#C0C0C0', 'x': 1,   'y': 13.0, 'w': 7, 'h': 2.5},
    'F1': {'text': 'Breaking alignment into\nk-mers', 'shape': 'box', 'fillcolor': '#C0C0C0', 'x': 1,   'y': 15.5, 'w': 7, 'h': 2.5},
    'F': {'text': 'Generate Word2Vec vectors for each k-mer,\naggregate them by mean to form a single sentence vector per sequence', 'shape': 'box', 'fillcolor': '#B3B3B3', 'x': 1, 'y': 18.0, 'w': 7, 'h': 2.5},
    'F2': {'text': 'Save the global Word2Vec model', 'shape': 'box', 'fillcolor': '#B3B3B3', 'x': 1, 'y': 19.25, 'w': 7, 'h': 2.5},
    'G': {'text': 'Standardize Training Embeddings\n(StandardScaler)', 'shape': 'box', 'fillcolor': '#A6A6A6', 'x': 1,   'y': 20.5, 'w': 7, 'h': 2.5},

    # Painel B: Oversampling Phase (agora dois passos)
    'H':  {'text': 'Oversampling (Random oversampling and SMOTE)', 'shape': 'box', 'fillcolor': '#999999', 'x': 1,   'y': 23.0, 'w': 7, 'h': 2.5},
    'H2': {'text': 'Grid Search and Cross validation', 'shape': 'box', 'fillcolor': '#949494', 'x': 1,   'y': 25.0, 'w': 7, 'h': 2.5},

    # Painel C: Training and Optimization of Random Forest Phase (nÃ³s I e J)
    'I': {'text': 'Calibrate Model (Isotonic Calibration)', 'shape': 'box', 'fillcolor': '#8C8C8C', 'x': 1,   'y': 27.0, 'w': 7, 'h': 2.5},
    'J': {'text': 'Evaluate Model (ROC, F1, PR AUC)', 'shape': 'box', 'fillcolor': '#808080', 'x': 1,   'y': 29.0, 'w': 7, 'h': 2.5},

    # Painel D: Classification of New Sequences Phase â€“ Transfer Learning (nÃ³s Y, K, L, M, N, O, P, Z, T, Q, R)
    'Y': {'text': 'Extract FAAL Domain from\nPrediction FASTA (Auxiliary tool)', 'shape': 'box', 'fillcolor': '#7A7A7A', 'x': 1,   'y': 31.0, 'w': 7, 'h': 2.5},
    'K': {'text': 'Load Prediction Data (FASTA)', 'shape': 'parallelogram', 'fillcolor': '#737373', 'x': 1,   'y': 33.5, 'w': 7, 'h': 2.5},
    'L': {'text': 'Breaking New sequence into\nK-mers', 'shape': 'box', 'fillcolor': '#666666', 'x': 1,   'y': 36.0, 'w': 7, 'h': 2.5},
    'M': {'text': 'Generate New Sequence Embeddings\n(From k-mers)\nusing global ww model\n- (Transfer learning)', 'shape': 'box', 'fillcolor': '#5A5A5A', 'x': 1,   'y': 38.5, 'w': 7, 'h': 2.5},
    'N': {'text': 'Standardize New Sequence Embeddings\n(Using Training Scaler)', 'shape': 'box', 'fillcolor': '#4D4D4D', 'x': 1,   'y': 41.0, 'w': 7, 'h': 2.5},
    'O': {'text': 'Import Random Forest Classifier\nand RF Scaler', 'shape': 'box', 'fillcolor': '#404040', 'x': 1,   'y': 43.0, 'h': 2.5, 'w': 7},
    'P': {'text': 'Predict on New Sequences &amp;\nGet Class Rankings', 'shape': 'box', 'fillcolor': '#404040', 'x': 1,   'y': 45.0, 'w': 7, 'h': 2.5}, 		
    'Z': {'text': 'Sum Probabilities and Normalize', 'shape': 'box', 'fillcolor': '#505050', 'x': 1,   'y': 47.0, 'w': 7, 'h': 2.5},
    'T': {'text': 'Get Prediction Confidence', 'shape': 'box', 'fillcolor': '#414141', 'x': 1,   'y': 49.0, 'w': 7, 'h': 2.5},
    'Q': {'text': 'Visualize Results\n(Scatter Plot, Learning Curve,\nTables, UMAP)', 'shape': 'box', 'fillcolor': '#333333', 'x': 1,   'y': 51.0, 'w': 7, 'h': 2.5},
    'R': {'text': 'Output Results\n(Save Files, Download)', 'shape': 'ellipse', 'fillcolor': '#262626', 'x': 1,   'y': 53.5, 'w': 7, 'h': 2.5}
}

# ConexÃµes entre os nÃ³s (fluxo completo)
ppt_edges = [
    # Treinamento:
    ('A', 'X', ''),
    ('X', 'B', ''),
    ('B', 'C', ''),
    ('C', 'D', 'Not Aligned'),
    ('C', 'E', 'Aligned'),
    ('D', 'E', ''),
    ('E', 'F1', ''),
    ('F1', 'F', ''),
    ('F', 'F2', ''),
    ('F2', 'G', ''),
    ('G', 'H', ''),
    ('H', 'H2', ''),
    ('H2', 'I', ''),
    ('I', 'J', ''),
    # TransiÃ§Ã£o para prediÃ§Ã£o:
    ('J', 'Y', ''),
    # PrediÃ§Ã£o:
    ('Y', 'K', ''),
    ('K', 'L', ''),
    ('L', 'M', ''),
    ('M', 'N', ''),
    ('N', 'O', ''),
    ('O', 'P', ''),
    ('P', 'Z', ''),
    ('Z', 'T', ''),
    ('T', 'Q', ''),
    ('Q', 'R', '')
]

def create_ppt_flowchart(ppt_filename="workflow_native.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Slide em branco
    shapes_dict = {}

    # Adiciona cada nÃ³ com as propriedades definidas
    for node_id, props in ppt_nodes.items():
        shape_type = map_shape(props['shape'])
        left = Inches(props['x'])
        top = Inches(props['y'])
        width = Inches(props['w'])
        height = Inches(props['h'])
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(props['fillcolor'].lstrip('#'))
        shape.text = props['text']
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.name = 'Arial'
            paragraph.font.bold = True
            paragraph.alignment = 1  # Centralizado
            paragraph.font.color.rgb = RGBColor.from_string(ppt_font_color(props['fillcolor']))
        shapes_dict[node_id] = shape

    # Adiciona os tÃ­tulos dos painÃ©is como caixas de texto
    panel_titles = {
        'A': {'text': 'A) Neural Network: Protein Embedding Phase', 'x': Inches(1), 'y': Inches(0),   'w': Inches(7), 'h': Inches(0.8)},
        'B': {'text': 'B) Oversampling Phase',                     'x': Inches(1), 'y': Inches(21.0), 'w': Inches(7), 'h': Inches(0.8)},
        'C': {'text': 'C) Training and Optimization of Random Forest Phase', 'x': Inches(1), 'y': Inches(26.5), 'w': Inches(7), 'h': Inches(0.8)},
        'D': {'text': 'D) Classification of New Sequences Phase â€“ Transfer Learning', 'x': Inches(1), 'y': Inches(31.0), 'w': Inches(7), 'h': Inches(0.8)}
    }
    for p in panel_titles.values():
        title_box = slide.shapes.add_textbox(p['x'], p['y'], p['w'], p['h'])
        title_tf = title_box.text_frame
        title_tf.text = p['text']
        for para in title_tf.paragraphs:
            para.font.size = Pt(24)
            para.font.name = 'Arial'
            para.font.bold = True
            para.alignment = 1

    # Adiciona os conectores entre os nÃ³s
    for start_id, end_id, label in ppt_edges:
        start_shape = shapes_dict[start_id]
        end_shape = shapes_dict[end_id]
        start_x = int(start_shape.left + start_shape.width)
        start_y = int(start_shape.top + start_shape.height / 2)
        end_x = int(end_shape.left)
        end_y = int(end_shape.top + end_shape.height / 2)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, start_y,
            end_x - start_x,
            end_y - start_y
        )
        connector.line.width = Pt(2)
        connector.begin_connect(start_shape, 1)
        connector.end_connect(end_shape, 3)
        if label:
            mid_x = int((start_x + end_x) / 2)
            mid_y = int((start_y + end_y) / 2)
            textbox = slide.shapes.add_textbox(Inches(mid_x/72), Inches(mid_y/72), Inches(1), Inches(0.5))
            textbox.text_frame.text = label
            for paragraph in textbox.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.name = 'Arial'
                paragraph.font.bold = True

    prs.save(ppt_filename)
    print(f"PPT salvo como {ppt_filename}")

# ============================================================
# EXECUÃ‡ÃƒO: GERANDO OS FORMATOS DO FLUXOGRAMA
# ============================================================
if __name__ == "__main__":
    # Gera imagens PNG, SVG e PDF com DPI 900
    generate_flowchart_images("workflow")
    # Gera arquivo PPTX com os painÃ©is organizados verticalmente, incluindo as etapas divididas de Word2Vec.
    create_ppt_flowchart("workflow_native.pptx")
    
# ====================================================================
# (O restante do cÃ³digo, que inclui funÃ§Ãµes de processamento, treinamento, visualizaÃ§Ãµes e interface Streamlit,
# permanece inalterado, exceto pelas modificaÃ§Ãµes realizadas na parte do fluxograma.)
# ====================================================================

